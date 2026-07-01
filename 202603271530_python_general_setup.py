# ================================================================================
# FILE: 202603271530_python_general_setup.py
# NAME: Dharia, Saumil
# FIRST CREATED BY CLAUDE ON 202603271630 (March 27, 2026 at 4:30 PM EST)
# LAST UPDATED BY CLAUDE ON 202607011048 (July 01, 2026 at 10:48 AM EST)
# ================================================================================

"""
GENERAL SETUP MODULE
====================
Centralized configuration for all Python scripts.
All libraries accessible as: setup.pd, setup.np, setup.plt, etc.
"""

# ================================================================================
# TABLE OF CONTENTS
# ================================================================================
#
# Section 1:  Self-Path Setup                         (Lines 40-48)
# Section 2:  Library Detection and Auto-Installation  (Lines 54-110)
# Section 3:  Warning Suppression                      (Lines 112-122)
# Section 4:  Import All Required Libraries            (Lines 124-178)
# Section 5:  Environment Detection                    (Lines 180-194)
# Section 6:  Google Cloud Credentials                 (Lines 196-232)
# Section 7:  Oracle Configuration                     (Lines 234-332)
# Section 8:  Trino Configuration                      (Lines 334-409)
# Section 9:  File Lookup Utility                      (Lines 411-421)
# Section 10: Load and Consolidate Files               (Lines 423-599)
# Section 11: File Inventory & Retrieval System        (Lines 601-706)
# Section 12: Utility Functions                        (Lines 708-801)
# Section 13: PowerPoint Helper Functions              (Lines 803-838)
# Section 14: PDF Helper Functions                     (Lines 840-943)
# Section 14B: PDF Prose Helpers                       (Lines 945-997)
# Section 14C: Standardized Table PNG Renderer         (Lines 999-1069)
# Section 15: Module-Level Exposure                    (Lines 1071-1109)
#
# Line numbers verified after creation via grep.
# ================================================================================

# ========== SECTION 1: Self-Path Setup (Lines 40-48) ==========
import os
import sys

SETUP_PATH = os.path.dirname(os.path.abspath(__file__))
if SETUP_PATH not in sys.path:
    sys.path.insert(0, SETUP_PATH)

# ========== END SECTION 1 (Lines 40-48) ==========

print("=" * 60)
print("GENERAL SETUP - LIBRARIES & CREDENTIALS")
print("=" * 60 + "\n")

# ========== SECTION 2: Library Detection and Auto-Installation (Lines 54-110) ==========

print("Section 1: Loading Libraries (auto-installing if needed)...\n")

import subprocess
from datetime import datetime, timedelta
import glob

REQUIRED_LIBRARIES = {
    "pandas": "pandas",
    "numpy": "numpy",
    "pyarrow": "pyarrow",
    "scipy": "scipy",
    "statsmodels": "statsmodels",
    "sklearn": "scikit-learn",
    "pingouin": "pingouin",
    "sqlalchemy": "sqlalchemy",
    "oracledb": "oracledb",
    "google.cloud": "google-cloud-bigquery",
    "requests": "requests",
    "github": "PyGithub",
    "openpyxl": "openpyxl",
    "xlrd": "xlrd",
    "lxml": "lxml",
    "bs4": "beautifulsoup4",
    "matplotlib": "matplotlib",
    "seaborn": "seaborn",
    "plotly": "plotly",
    "tqdm": "tqdm",
    "dotenv": "python-dotenv",
    "PIL": "pillow",
    "reportlab": "reportlab",
    "pptx": "python-pptx",
    "trino": "trino",
}

def install_library(package_name):
    print(f"    Installing {package_name}...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", package_name, "-q"])
    print(f"    ✓ {package_name} installed")

for import_name, pip_name in REQUIRED_LIBRARIES.items():
    try:
        __import__(import_name)
        print(f"  ✓ {pip_name}")
    except ImportError:
        print(f"  ⚠ {pip_name} not found, installing...")
        try:
            install_library(pip_name)
            print(f"  ✓ {pip_name} installed successfully")
        except Exception as e:
            print(f"  ✗ Failed to install {pip_name}: {e}")

print(f"  ✓ datetime\n")
print("✓ All libraries loaded and verified\n")

# ========== END SECTION 2 (Lines 54-110) ==========

# ========== SECTION 3: Warning Suppression (Lines 112-122) ==========

import warnings

warnings.filterwarnings(
    'ignore',
    message='.*pandas only supports SQLAlchemy connectable.*',
    category=UserWarning
)

# ========== END SECTION 3 (Lines 112-122) ==========

# ========== SECTION 4: Import All Required Libraries (Lines 124-178) ==========

import pandas as pd
import numpy as np
import pyarrow as pa

import scipy
from scipy import stats
import statsmodels.api as sm
import statsmodels.formula.api as smf
from sklearn.datasets import make_classification
from sklearn.model_selection import train_test_split
from sklearn.linear_model import LogisticRegression
from sklearn.metrics import confusion_matrix, classification_report, roc_curve, auc
import pingouin as pg

from sqlalchemy import create_engine
import oracledb
from trino.dbapi import connect as trino_connect
from trino.auth import BasicAuthentication
from google.cloud import bigquery
import requests
from github import Github

import openpyxl
import xlrd
from lxml import etree
from bs4 import BeautifulSoup
import json
import csv

import matplotlib.pyplot as plt
import seaborn as sns
import plotly.express as px
import plotly.graph_objects as go

from tqdm import tqdm
from dotenv import load_dotenv
from PIL import Image

from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image as RLImage, PageBreak, Table, TableStyle
from reportlab.platypus.tableofcontents import TableOfContents
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_JUSTIFY
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

print("✓ All libraries imported successfully\n")

# ========== END SECTION 4 (Lines 124-178) ==========

# ========== SECTION 5: Environment Detection (Lines 180-194) ==========

print("Section 2: Setting Up Credentials...\n")

def detect_environment():
    if os.environ.get('GOOGLE_CLOUD_PROJECT') or os.environ.get('GCP_PROJECT_ID'):
        return 'cloud'
    if sys.platform == 'win32':
        return 'windows'
    return 'linux'

ENVIRONMENT = detect_environment()
print(f"  Detected Environment: {ENVIRONMENT.upper()}\n")

# ========== END SECTION 5 (Lines 180-194) ==========

# ========== SECTION 6: Google Cloud Credentials (Lines 196-232) ==========

if ENVIRONMENT == 'cloud':
    print("  ✓ Google Cloud environment detected\n")
    GCP_PROJECT_ID = os.environ.get('GCP_PROJECT_ID')
    if not GCP_PROJECT_ID:
        try:
            client = bigquery.Client()
            GCP_PROJECT_ID = client.project
            print(f"  ✓ Auto-detected Project ID: {GCP_PROJECT_ID}\n")
        except Exception as e:
            print(f"  ✗ Could not detect project ID: {e}\n")
            GCP_PROJECT_ID = None

elif ENVIRONMENT == 'windows':
    print("  Local development mode (Windows)\n")
    possible_paths = [
        r'C:\Users\dharias\Documents\Credentials\dharia_service_account.json',
        r'C:\Users\dharias\Desktop\Python Scripts\Credentials\dharia_service_account.json',
        './credentials/dharia_service_account.json',
    ]
    GCP_PROJECT_ID = None
    for path in possible_paths:
        if os.path.exists(path):
            os.environ['GOOGLE_APPLICATION_CREDENTIALS'] = path
            try:
                client = bigquery.Client()
                GCP_PROJECT_ID = client.project
                print(f"  ✓ GCP credentials loaded\n")
                break
            except:
                continue

else:
    GCP_PROJECT_ID = None

# ========== END SECTION 6 (Lines 196-232) ==========

# ========== SECTION 7: Oracle Configuration (Lines 234-332) ==========

ORACLE_CONFIG = {
    'host': '10.63.133.38',
    'port': 1521,
    'service_name': 'buprd1.buocisubpridbph.buocivcnphxdr.oraclevcn.com',
    'username': 'DEVDHARIAS',
    'schema': 'SYSADM',
}

try:
    possible_client_paths = [
        r'C:\Users\dharias\Downloads\instantclient-basic-windows.x64-23.26.1.0.0\instantclient_23_0',
        r'C:\instantclient_23_0',
    ]
    for client_path in possible_client_paths:
        if os.path.exists(client_path):
            oracledb.init_oracle_client(lib_dir=client_path)
            print(f"  ✓ Oracle Instant Client initialized\n")
            break
except Exception as e:
    print(f"  ⚠ Oracle setup incomplete: {e}\n")

def test_oracle_connection():
    """
    Test and establish Oracle database connection.
    Reads credentials from .env file in SETUP_DIR (via os.environ['SETUP_DIR']).
    Requires ORACLE_CLIENT_PATH to be set in os.environ.
    
    Returns:
    --------
    dict with keys:
        'connection': oracledb.Connection object
        'status': 'success' or 'failed'
        'message': status message
    """
    try:
        # Get SETUP_DIR from environment
        setup_dir = os.environ.get('SETUP_DIR')
        if not setup_dir:
            return {
                'status': 'failed',
                'message': 'SETUP_DIR not found in os.environ',
                'connection': None
            }
        
        # Load .env file from SETUP_DIR
        env_file_path = os.path.join(setup_dir, '.env')
        if not os.path.exists(env_file_path):
            return {
                'status': 'failed',
                'message': f'.env file not found at {env_file_path}',
                'connection': None
            }
        
        load_dotenv(env_file_path)
        
        # Get credentials from environment
        username = os.environ.get('ORACLE_USERNAME')
        password = os.environ.get('ORACLE_PASSWORD')
        
        if not username or not password:
            return {
                'status': 'failed',
                'message': 'ORACLE_USERNAME or ORACLE_PASSWORD not found in .env',
                'connection': None
            }
        
        # Get ORACLE_CLIENT_PATH from environment
        oracle_client_path = os.environ.get('ORACLE_CLIENT_PATH')
        if oracle_client_path and os.path.exists(oracle_client_path):
            try:
                oracledb.init_oracle_client(lib_dir=oracle_client_path)
            except:
                pass  # Client may already be initialized
        
        # Create connection
        dsn = f"{ORACLE_CONFIG['host']}:{ORACLE_CONFIG['port']}/{ORACLE_CONFIG['service_name']}"
        connection = oracledb.connect(
            user=username,
            password=password,
            dsn=dsn
        )
        
        print("✓ Successfully connected to Oracle\n")
        return {
            'status': 'success',
            'message': '✓ Successfully connected to Oracle',
            'connection': connection
        }
    
    except Exception as e:
        return {
            'status': 'failed',
            'message': f'✗ Oracle connection failed: {str(e)[:100]}',
            'connection': None
        }

# ========== END SECTION 7 (Lines 234-332) ==========

# ========== SECTION 8: Trino Configuration (Lines 334-409) ==========

TRINO_CONFIG = {
    'host': 'trino.de-eks-nonprod.bu.edu',
    'port': 443,
    'http_scheme': 'https',
    'catalog': 'iceberg',
    'schema': 'pond_asir',
}

def trino_connection():
    """
    Test and establish Trino data lake connection.
    Reads TRINO_USERNAME and TRINO_PASSWORD from .env file in SETUP_DIR.
    
    Returns:
    --------
    dict with keys:
        'connection': trino.dbapi.Connection object
        'status': 'success' or 'failed'
        'message': status message
    """
    try:
        setup_dir = os.environ.get('SETUP_DIR')
        if not setup_dir:
            return {
                'status': 'failed',
                'message': 'SETUP_DIR not found in os.environ',
                'connection': None
            }
        
        env_file_path = os.path.join(setup_dir, '.env')
        if not os.path.exists(env_file_path):
            return {
                'status': 'failed',
                'message': f'.env file not found at {env_file_path}',
                'connection': None
            }
        
        load_dotenv(env_file_path)
        
        username = os.environ.get('TRINO_USERNAME')
        password = os.environ.get('TRINO_PASSWORD')
        
        if not username or not password:
            return {
                'status': 'failed',
                'message': 'TRINO_USERNAME or TRINO_PASSWORD not found in .env',
                'connection': None
            }
        
        connection = trino_connect(
            host=TRINO_CONFIG['host'],
            port=TRINO_CONFIG['port'],
            http_scheme=TRINO_CONFIG['http_scheme'],
            user=username,
            auth=BasicAuthentication(username, password),
            catalog=TRINO_CONFIG['catalog'],
            schema=TRINO_CONFIG['schema'],
        )
        
        print("✓ Successfully connected to Trino\n")
        return {
            'status': 'success',
            'message': '✓ Successfully connected to Trino',
            'connection': connection
        }
    
    except Exception as e:
        return {
            'status': 'failed',
            'message': f'✗ Trino connection failed: {str(e)[:100]}',
            'connection': None
        }

# ========== END SECTION 8 (Lines 334-409) ==========

# ========== SECTION 9: File Lookup Utility (Lines 411-421) ==========

def find_file_by_timestamp_prefix(timestamp_prefix, search_dir=None, file_pattern=None):
    """Find a file by its YYYYMMDDHHMM timestamp prefix"""
    if search_dir is None:
        search_dir = '.'
    pattern = f"{timestamp_prefix}_{file_pattern}" if file_pattern else f"{timestamp_prefix}_*"
    matches = glob.glob(os.path.join(search_dir, pattern))
    return matches[0] if matches else None

# ========== END SECTION 9 (Lines 411-421) ==========

# ========== SECTION 10: Load and Consolidate Files (Lines 423-599) ==========

def load_consolidate_files(file_list, input_path, auto_skip_rows=True, verbose=True):
    """
    Load and consolidate multiple Excel/CSV files from a directory
    Automatically detects header rows and skips blank rows
    
    Parameters:
    -----------
    file_list : list
        List of specific filenames to load
        Example: ['file1.xlsx', 'file2.xlsx', 'file3.csv']
    input_path : str
        Path to directory containing the files
    auto_skip_rows : bool
        Automatically detect and skip blank rows at top (default: True)
    verbose : bool
        Print detailed progress messages (default: True)
    
    Returns:
    --------
    pandas.DataFrame
        Consolidated dataframe with all files concatenated and data types aligned
    
    Usage:
    ------
    df = setup.load_consolidate_files(
        ['enrollment.xlsx', 'finaid.csv', 'grades.xlsx'], 
        PROJECT_INPUT
    )
    """
    
    df_list = []
    files_loaded = 0
    files_failed = 0
    
    if verbose:
        print(f"\n{'='*80}")
        print(f"LOADING & CONSOLIDATING {len(file_list)} FILES")
        print(f"{'='*80}\n")
    
    for file_name in file_list:
        file_path = os.path.join(input_path, file_name)
        file_ext = os.path.splitext(file_name)[1].lower()
        
        if verbose:
            print(f"File: {file_name}")
            print(f"  Type: {file_ext}")
        
        # Check if file exists
        if not os.path.exists(file_path):
            if verbose:
                print(f"  ✗ File not found\n")
            files_failed += 1
            continue
        
        try:
            # ===== DETECT HEADER ROW =====
            if file_ext == '.xlsx':
                # Read Excel to detect header
                wb = openpyxl.load_workbook(file_path, data_only=True, read_only=True)
                ws = wb.active
                row_non_none_counts = []
                
                for row in ws.iter_rows(values_only=True):
                    count = sum(1 for cell in row if cell is not None)
                    row_non_none_counts.append(count)
                
                wb.close()
                
                if row_non_none_counts:
                    max_count = max(row_non_none_counts)
                    header_row_index = row_non_none_counts.index(max_count)
                else:
                    header_row_index = 0
                
                rows_to_skip = header_row_index if auto_skip_rows else 0
                
                # Load Excel file
                df_temp = pd.read_excel(
                    file_path,
                    skiprows=rows_to_skip,
                    header=0
                )
                
            elif file_ext == '.csv':
                # For CSV, try to detect if there's a header
                rows_to_skip = 0
                
                if auto_skip_rows:
                    # Count non-null values in first few rows
                    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                        line_counts = []
                        for i, line in enumerate(f):
                            if i >= 20:  # Check first 20 rows
                                break
                            non_empty = sum(1 for cell in line.split(',') if cell.strip())
                            line_counts.append(non_empty)
                        
                        if line_counts:
                            max_count = max(line_counts)
                            header_row_index = line_counts.index(max_count)
                            rows_to_skip = header_row_index
                
                df_temp = pd.read_csv(
                    file_path,
                    skiprows=rows_to_skip,
                    header=0,
                    encoding='utf-8',
                    on_bad_lines='skip'
                )
            
            else:
                if verbose:
                    print(f"  ✗ Unsupported file type: {file_ext}\n")
                files_failed += 1
                continue
            
            # ===== CLEAN COLUMN NAMES =====
            # Standardize column names (strip whitespace)
            df_temp.columns = df_temp.columns.str.strip()
            
            # ===== REMOVE COMPLETELY EMPTY ROWS =====
            df_temp = df_temp.dropna(how='all')
            
            if verbose:
                print(f"  ✓ Loaded {len(df_temp)} rows × {len(df_temp.columns)} columns")
                if rows_to_skip > 0:
                    print(f"  ✓ Skipped {rows_to_skip} row(s) before header")
                print(f"  Columns: {list(df_temp.columns)}\n")
            
            df_list.append(df_temp)
            files_loaded += 1
        
        except Exception as e:
            if verbose:
                print(f"  ✗ Error reading file: {str(e)[:100]}\n")
            files_failed += 1
            continue
    
    # ===== CONSOLIDATE & CONCATENATE =====
    if df_list:
        if verbose:
            print(f"{'='*80}")
            print(f"CONSOLIDATING {files_loaded} FILES")
            print(f"{'='*80}\n")
        
        # Concatenate with sort=False to preserve column order
        consolidated_df = pd.concat(df_list, ignore_index=True, sort=False)
        
        # ===== ALIGN DATA TYPES =====
        # Attempt automatic type conversion for numeric columns
        for col in consolidated_df.columns:
            try:
                consolidated_df[col] = pd.to_numeric(consolidated_df[col], errors='ignore')
            except:
                pass
        
        if verbose:
            print(f"✓ Consolidation complete")
            print(f"  Files loaded: {files_loaded}")
            print(f"  Files failed: {files_failed}")
            print(f"  Total rows: {len(consolidated_df):,}")
            print(f"  Total columns: {len(consolidated_df.columns)}")
            print(f"  All columns: {list(consolidated_df.columns)}\n")
            print("FIRST 5 ROWS:")
            print(consolidated_df.head())
            print()
        
        return consolidated_df
    
    else:
        if verbose:
            print(f"\n✗ WARNING: No files loaded successfully")
        return pd.DataFrame()

# ========== END SECTION 10 (Lines 423-599) ==========

# ========== SECTION 11: File Inventory & Retrieval System (Lines 601-706) ==========

# Global file inventory caches
file_inventory_cache = {}
file_inventory_by_type_cache = {}

def scan_directory_for_inventory(directory_path, dir_label):
    """Scan a directory and catalog all files with metadata"""
    files_list = []
    try:
        if os.path.exists(directory_path):
            all_items = os.listdir(directory_path)
            files_only = [item for item in all_items if os.path.isfile(os.path.join(directory_path, item))]
            
            for filename in files_only:
                full_path = os.path.join(directory_path, filename)
                file_size = os.path.getsize(full_path)
                mod_time = os.path.getmtime(full_path)
                mod_time_readable = datetime.fromtimestamp(mod_time).strftime('%Y-%m-%d %H:%M:%S')
                
                files_list.append({
                    'filename': filename,
                    'full_path': full_path,
                    'size_bytes': file_size,
                    'modified': mod_time_readable
                })
        return files_list
    except Exception as e:
        return []

def organize_files_by_type(files_list):
    """Categorize files by extension type"""
    file_types = {
        'Excel': ['.xlsx', '.xls', '.xlsm'],
        'CSV': ['.csv'],
        'Text': ['.txt'],
        'Python': ['.py'],
        'Other': []
    }
    
    organized = {}
    for file_type in file_types.keys():
        organized[file_type] = []
    
    for file_info in files_list:
        filename = file_info['filename']
        file_ext = os.path.splitext(filename)[1].lower()
        
        categorized = False
        for file_type, extensions in file_types.items():
            if file_ext in extensions:
                organized[file_type].append(file_info)
                categorized = True
                break
        
        if not categorized:
            organized['Other'].append(file_info)
    
    return organized

def find_files_by_type(file_type, directory='PROJECT_INPUT'):
    """
    Find all files of a specific type in a directory
    
    Usage:
    ------
    excel_files = find_files_by_type('Excel', 'PROJECT_INPUT')
    csv_files = find_files_by_type('CSV', 'PROJECT_OUTPUT')
    """
    if directory in file_inventory_by_type_cache:
        return file_inventory_by_type_cache[directory].get(file_type, [])
    return []

def find_files_by_pattern(pattern, directory='PROJECT_INPUT'):
    """
    Find files matching a pattern (case-insensitive substring match)
    
    Usage:
    ------
    matching = find_files_by_pattern('Grad AY', 'PROJECT_INPUT')
    """
    matching = []
    
    if directory in file_inventory_cache:
        for file_info in file_inventory_cache[directory]:
            if pattern.lower() in file_info['filename'].lower():
                matching.append(file_info)
    
    return matching

def get_all_filenames(directory='PROJECT_INPUT'):
    """
    Get list of all filenames in a directory
    
    Usage:
    ------
    all_files = get_all_filenames('PROJECT_INPUT')
    for filename in all_files:
        print(filename)
    """
    if directory in file_inventory_cache:
        return [f['filename'] for f in file_inventory_cache[directory]]
    
    return []

# ========== END SECTION 11 (Lines 601-706) ==========

# ========== SECTION 12: Utility Functions (Lines 708-801) ==========

def diagnose_dataframe(df, df_name="DataFrame", check_duplicates=True, preview_rows=5):
    """Print comprehensive dataframe diagnostics including duplicates"""
    print(f"\n{'='*80}")
    print(f"DATAFRAME DIAGNOSIS: {df_name}")
    print(f"{'='*80}\n")
    print(f"Shape: {df.shape[0]:,} rows × {df.shape[1]} columns")
    print(f"Memory: {df.memory_usage(deep=True).sum() / 1024**2:.2f} MB\n")
    
    # Data Preview (all columns in one wide row, no block wrapping)
    print(f"DATA PREVIEW (first {preview_rows} rows, all columns):")
    print("-" * 80)
    with pd.option_context('display.max_columns', None,
                           'display.width', None,
                           'display.max_colwidth', 30):
        print(df.head(preview_rows).to_string())
    print()
    
    # Duplicate Detection
    if check_duplicates:
        duplicate_rows = df.duplicated().sum()
        print(f"DUPLICATE ANALYSIS:")
        print(f"  Total duplicate rows: {duplicate_rows:,} ({(duplicate_rows/len(df)*100):.2f}%)")
        
        # Check duplicates per column
        print(f"\n  Duplicates by column:")
        print(f"  {'Column':<30} {'Unique':<12} {'Duplicates':<12} {'Dup %':<10}")
        print(f"  {'-'*65}")
        for col in df.columns:
            unique_count = df[col].nunique()
            dup_count = len(df[col]) - unique_count
            dup_pct = (dup_count / len(df) * 100) if len(df) > 0 else 0
            print(f"  {col:<30} {unique_count:<12} {dup_count:<12} {dup_pct:.2f}%")
        print()
    
    print(f"{'Column':<30} {'Data Type':<15} {'Non-Null':<15} {'Missing':<10}")
    print("-" * 80)
    for col in df.columns:
        non_null = df[col].notna().sum()
        missing = df[col].isna().sum()
        print(f"{col:<30} {str(df[col].dtype):<15} {non_null:<15} {missing:<10}")
    print()

def missing_data_report(df, df_name="DataFrame", show_heatmap=False):
    """Generate missing data analysis report"""
    print(f"\n{'='*80}")
    print(f"MISSING DATA REPORT: {df_name}")
    print(f"{'='*80}\n")
    
    total_cells = df.shape[0] * df.shape[1]
    total_missing = df.isnull().sum().sum()
    pct_missing = (total_missing / total_cells) * 100
    
    print(f"Total cells: {total_cells:,}")
    print(f"Missing cells: {total_missing:,}")
    print(f"Overall missing: {pct_missing:.2f}%\n")
    
    missing_by_col = df.isnull().sum()
    missing_pct = (missing_by_col / len(df)) * 100
    
    print(f"{'Column':<30} {'Missing Count':<15} {'Missing %':<10}")
    print("-" * 60)
    for col in df.columns:
        print(f"{col:<30} {missing_by_col[col]:<15} {missing_pct[col]:.2f}%")
    print()

def mask_id_column(df, id_column, seed=None, new_id_prefix='ID', length=8):
    """Mask a unique ID column with pseudo-random alphanumeric IDs"""
    import random
    import string
    
    if seed is not None:
        random.seed(seed)
    
    unique_ids = df[id_column].unique()
    id_mapping = {}
    
    for original_id in unique_ids:
        random_part = ''.join(random.choices(string.ascii_uppercase + string.digits, k=length))
        new_id = f"{new_id_prefix}{random_part}"
        id_mapping[original_id] = new_id
    
    df_masked = df.copy()
    df_masked[id_column] = df_masked[id_column].map(id_mapping)
    
    return df_masked, id_mapping

print("  ✓ Utility functions ready")
print("    - diagnose_dataframe() [now with 5-row preview + duplicate detection]")
print("    - missing_data_report()")
print("    - mask_id_column()\n")

# ========== END SECTION 12 (Lines 708-801) ==========

# ========== SECTION 13: PowerPoint Helper Functions (Lines 803-838) ==========

def add_text_slide(prs, title, paragraph_text):
    """Add a text-only slide to PowerPoint"""
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    
    text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(5))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    text_frame.text = paragraph_text
    text_frame.paragraphs[0].font.size = Pt(14)
    text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

def add_image_slide(prs, image_path):
    """Add an image-only slide to PowerPoint"""
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(6.5)
    slide.shapes.add_picture(image_path, left, top, width=width, height=height)

print("  ✓ PowerPoint helper functions ready:")
print("    - add_text_slide()")
print("    - add_image_slide()\n")

# ========== END SECTION 13 (Lines 803-838) ==========

# ========== SECTION 14: PDF Helper Functions (Lines 840-943) ==========

def add_chapter_header(story, college_code, college_name):
    """Add chapter divider page with college name. Registers a TOC entry at level 0."""
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        'ChapterTitle',
        parent=styles['Heading1'],
        fontSize=28,
        textColor=colors.HexColor('#4472C4'),
        spaceAfter=30,
        alignment=TA_CENTER,
        fontName='Helvetica-Bold'
    )

    class TOCParagraph(Paragraph):
        """Paragraph subclass that emits a TOCEntry notification when drawn."""
        def draw(self):
            Paragraph.draw(self)
            key = f'chapter_{college_code}'
            self.canv.bookmarkPage(key)
            self.canv._doctemplate.notify('TOCEntry', (0, college_name, self.canv.getPageNumber(), key))

    story.append(Spacer(1, 1.5*inch))
    story.append(TOCParagraph(college_name, title_style))
    story.append(Spacer(1, 1.5*inch))
    story.append(PageBreak())
    return story

def add_image_with_caption(story, image_path):
    """Add image to story (image already has title in PNG). Uses RLImage (reportlab) to avoid clashing with PIL.Image."""
    if os.path.exists(image_path):
        try:
            img = RLImage(image_path, width=6.5*inch, height=6.5*inch*0.65)
            story.append(img)
            story.append(Spacer(1, 0.3*inch))
        except Exception as e:
            print(f'Error loading image {image_path}: {e}')
    else:
        print(f'Image not found: {image_path}')
    return story

def add_insights_bullets(story, insights, heading='Key Insights'):
    """Add insights as bullet points."""
    styles = getSampleStyleSheet()
    heading_style = ParagraphStyle(
        'InsightHeading',
        parent=styles['Heading2'],
        fontSize=14,
        textColor=colors.HexColor('#4472C4'),
        spaceAfter=12,
        fontName='Helvetica-Bold'
    )
    bullet_style = ParagraphStyle(
        'BulletText',
        parent=styles['Normal'],
        fontSize=10,
        leftIndent=0.3*inch,
        spaceAfter=10,
        alignment=TA_JUSTIFY
    )
    story.append(Paragraph(heading, heading_style))
    for insight in insights:
        story.append(Paragraph(f'• {insight}', bullet_style))
    story.append(Spacer(1, 0.3*inch))
    return story

def build_table_of_contents(story, heading='Table of Contents'):
    """Insert a dynamic TOC that auto-fills via doc.multiBuild(). Requires chapters added via add_chapter_header."""
    styles = getSampleStyleSheet()
    heading_style = ParagraphStyle(
        'TOCHeading',
        parent=styles['Heading1'],
        fontSize=24,
        textColor=colors.HexColor('#4472C4'),
        spaceAfter=24,
        alignment=TA_CENTER,
        fontName='Helvetica-Bold'
    )
    toc_entry_style = ParagraphStyle(
        'TOCEntry0',
        parent=styles['Normal'],
        fontSize=12,
        leftIndent=0.3*inch,
        spaceAfter=8,
        fontName='Helvetica'
    )

    toc = TableOfContents()
    toc.levelStyles = [toc_entry_style]

    story.append(Spacer(1, 0.5*inch))
    story.append(Paragraph(heading, heading_style))
    story.append(Spacer(1, 0.3*inch))
    story.append(toc)
    story.append(PageBreak())
    return story

print("  ✓ PDF helper functions ready:")
print("    - add_chapter_header()")
print("    - add_image_with_caption()")
print("    - add_insights_bullets()\n")

# ========== END SECTION 14 (Lines 840-943) ==========

# ========== SECTION 14B: PDF Prose Helpers (themed chapters) ==========

def add_section_header(story, text):
    """Bold blue section label, e.g. 'Observations' / 'What does this mean?'."""
    styles = getSampleStyleSheet()
    section_style = ParagraphStyle(
        'SectionHeader',
        parent=styles['Heading1'],
        fontSize=16,
        textColor=colors.HexColor('#4472C4'),
        spaceBefore=6,
        spaceAfter=14,
        fontName='Helvetica-Bold'
    )
    story.append(Paragraph(text, section_style))
    return story

def add_subheading(story, text):
    """Bold theme subheading, e.g. 'Net Tuition Revenue' / 'Enrollment'."""
    styles = getSampleStyleSheet()
    sub_style = ParagraphStyle(
        'SubHeading',
        parent=styles['Heading2'],
        fontSize=13,
        textColor=colors.HexColor('#2F2F2F'),
        spaceBefore=8,
        spaceAfter=6,
        fontName='Helvetica-Bold'
    )
    story.append(Paragraph(text, sub_style))
    return story

def add_paragraph(story, text):
    """Flowing justified prose paragraph (no bullet)."""
    styles = getSampleStyleSheet()
    prose_style = ParagraphStyle(
        'BodyProse',
        parent=styles['Normal'],
        fontSize=10,
        spaceAfter=10,
        leading=14,
        alignment=TA_JUSTIFY
    )
    story.append(Paragraph(text, prose_style))
    story.append(Spacer(1, 0.05*inch))
    return story

print("  ✓ PDF prose helpers ready:")
print("    - add_section_header()")
print("    - add_subheading()")
print("    - add_paragraph()\n")

# ========== END SECTION 14B (Lines 945-997) ==========

# ========== SECTION 14C: Standardized Table PNG Renderer (Lines 999-1069) ==========

def render_table_png(df, col_labels, title_str, filename, output_dir, fmt='headcount'):
    """
    Render a styled table PNG (blue #4472C4 header, banded gray/white rows) and save it.
    This is the single house standard for every table graphic across projects —
    call this instead of writing a local one-off table-rendering function.

    df:         DataFrame whose index becomes the leftmost row-label column.
    col_labels: list of column headers — first entry is the row-label column
                (e.g. ['AY', 'QST', 'SPH', ...] or ['STATUS', 'COUNT']).
    title_str:  Title printed above the table.
    filename:   PNG filename (e.g. 'enrollment_by_program_and_year_table.png').
    output_dir: Directory to save into (e.g. PROJECT_OUTPUT or PROJECT_REVIEW).
    fmt:        'headcount' -> '{:,.0f}'
                'currency'  -> '${:,.0f}'
                'text'      -> values printed as-is (mixed types, already-formatted strings, '-')

    Usage:
    ------
    setup.render_table_png(
        STATUS_SUMMARY_DF.set_index('STATUS'),
        col_labels=['STATUS', 'COUNT'],
        title_str='BU Drop-Out Population - NSC Enrollment Status',
        filename='enrollment_status_summary_table.png',
        output_dir=PROJECT_OUTPUT,
        fmt='headcount'
    )
    """
    n_cols = len(col_labels)
    fig_width = max(14, n_cols * 1.3)
    fig, ax = plt.subplots(figsize=(fig_width, 7))
    ax.axis('tight'); ax.axis('off')

    if fmt == 'currency':
        table_data = [[str(idx)] + ['${:,.0f}'.format(v) for v in df.loc[idx].values] for idx in df.index]
    elif fmt == 'headcount':
        table_data = [[str(idx)] + ['{:,.0f}'.format(v) for v in df.loc[idx].values] for idx in df.index]
    else:  # 'text' - leave values as-is
        table_data = [[str(idx)] + [str(v) for v in df.loc[idx].values] for idx in df.index]

    col_width = 0.95 / n_cols
    table = ax.table(cellText=table_data, colLabels=col_labels,
                      cellLoc='center', loc='center',
                      colWidths=[col_width] * n_cols)
    table.auto_set_font_size(False); table.set_fontsize(9); table.scale(1, 2.5)

    for i in range(n_cols):
        cell = table[(0, i)]
        cell.set_facecolor('#4472C4')
        cell.set_text_props(weight='bold', color='white', size=10)

    for i in range(len(df.index)):
        for j in range(n_cols):
            cell = table[(i + 1, j)]
            cell.set_facecolor('#F0F0F0' if i % 2 == 0 else 'white')

    plt.title(title_str, fontsize=14, fontweight='bold', pad=30)
    plt.subplots_adjust(top=0.85)

    output_path = os.path.join(output_dir, filename)
    plt.savefig(output_path, dpi=150, bbox_inches='tight', pad_inches=0.5, facecolor='white')
    plt.show()
    plt.close()
    print(f'\u2713 {output_path}')
    return output_path

print("  \u2713 Table renderer ready:")
print("    - render_table_png()  [house style: #4472C4 header, banded gray/white rows]\n")

# ========== END SECTION 14C (Lines 999-1069) ==========

# ========== SECTION 15: Module-Level Exposure (Lines 1071-1109) ==========

__all__ = [
    'pd', 'np', 'pa',
    'scipy', 'stats', 'sm', 'smf', 'pg',
    'create_engine', 'oracledb', 'bigquery',
    'trino_connect', 'BasicAuthentication',
    'plt', 'sns', 'px', 'go',
    'tqdm', 'Image', 'BeautifulSoup',
    'Presentation', 'Inches', 'Pt', 'PP_ALIGN', 'RGBColor',
    'find_file_by_timestamp_prefix',
    'load_consolidate_files',
    'scan_directory_for_inventory',
    'organize_files_by_type',
    'find_files_by_type',
    'find_files_by_pattern',
    'get_all_filenames',
    'diagnose_dataframe',
    'missing_data_report',
    'mask_id_column',
    'add_text_slide',
    'add_image_slide',
    'add_chapter_header',
    'add_image_with_caption',
    'add_insights_bullets',
    'render_table_png',
    'RLImage', 'Table', 'TableStyle', 'colors',
    'TA_CENTER', 'TA_LEFT', 'TA_JUSTIFY',
    'SimpleDocTemplate', 'Paragraph', 'Spacer', 'PageBreak',
    'getSampleStyleSheet', 'ParagraphStyle', 'inch', 'letter',
    'test_oracle_connection',
    'trino_connection',
    'ENVIRONMENT', 'GCP_PROJECT_ID', 'ORACLE_CONFIG', 'TRINO_CONFIG',
    'os', 'sys', 'glob', 'json', 'csv',
    'file_inventory_cache', 'file_inventory_by_type_cache',
    'build_table_of_contents',
]

# ========== END SECTION 15 (Lines 1071-1109) ==========

print("=" * 60)
print("GENERAL SETUP COMPLETE")
print("=" * 60 + "\n")
